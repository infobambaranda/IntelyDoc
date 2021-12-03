#scanned_pdf and image
from pdf2image import convert_from_path
import pytesseract 
from PIL import Image
import os

#read_pdf
import PyPDF2

#read_docx
from docx import Document

#read_xlsx
import pandas as pd

#read_pptx
from pptx import Presentation

#read_csv
import csv

def read_scanned_pdf(PDF_file):
    pages = convert_from_path(PDF_file, 500) 

    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'  

    text=''
    image_counter = 1 
    for page in pages:  
        filename = "uploads/page_"+str(image_counter)+".jpg"
        page.save(filename, 'JPEG') 
        image_counter = image_counter + 1
        text_in_page = str(pytesseract.image_to_string(Image.open(filename)))
        text=text+text_in_page
        os.remove(filename)      
    return text

def read_pdf(PDF_file):
    pdfFileObj = open(PDF_file, 'rb')
    pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
    pdfReader.numPages
    pageObj = pdfReader.getPage(0)
    text=pageObj.extractText()
    if text == '':
        text = read_scanned_pdf(PDF_file)
    pdfFileObj.close()
    return text

def read_docx(name):
    document = Document(name)
    text=''
    for para in document.paragraphs:
        text = text + para.text        
    return text

def read_xlsx(name):
    table = pd.read_excel(name)
    #convert dataframe to a string
    text=table.to_string()        
    return text

def read_pptx(name):    
    prs = Presentation(name)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = text + shape.text        
    return text

def read_csv(name):
    with open(name, "r") as f:
        reader = csv.reader(f)
        text=''
        for row in reader:
            for ele in row:
                if len(ele) != 0:
                    character=ele[0]
                    if character.isalpha():
                        text=text+ele+' '
    return text

def read_image(name):
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'  
    text_in_page = str(pytesseract.image_to_string(Image.open(name)))
    text=text_in_page

    return text

def read_txt(name):
    with open(name) as f:
        text = f.read()
    return text